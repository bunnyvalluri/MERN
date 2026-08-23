import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_internhub_presentation():
    prs = Presentation()
    # 16:9 widescreen layout (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    themes = [
        {"name": "Royal Blue", "primary": RGBColor(30, 58, 138),  "secondary": RGBColor(59, 130, 246),  "accent": RGBColor(239, 246, 255), "card_bg": RGBColor(255, 255, 255), "border": RGBColor(191, 219, 254), "text_dark": RGBColor(15, 23, 42), "text_muted": RGBColor(51, 65, 85), "tag_bg": RGBColor(219, 234, 254)},
        {"name": "Sky Blue",   "primary": RGBColor(2, 132, 199),  "secondary": RGBColor(14, 165, 233),  "accent": RGBColor(240, 249, 255), "card_bg": RGBColor(255, 255, 255), "border": RGBColor(186, 230, 253), "text_dark": RGBColor(15, 23, 42), "text_muted": RGBColor(51, 65, 85), "tag_bg": RGBColor(224, 242, 254)},
        {"name": "Teal",       "primary": RGBColor(13, 148, 136), "secondary": RGBColor(20, 184, 166),  "accent": RGBColor(240, 253, 250), "card_bg": RGBColor(255, 255, 255), "border": RGBColor(153, 246, 228), "text_dark": RGBColor(15, 23, 42), "text_muted": RGBColor(51, 65, 85), "tag_bg": RGBColor(204, 251, 241)},
        {"name": "Emerald",    "primary": RGBColor(5, 150, 105),  "secondary": RGBColor(16, 185, 129),  "accent": RGBColor(236, 253, 245), "card_bg": RGBColor(255, 255, 255), "border": RGBColor(167, 243, 208), "text_dark": RGBColor(15, 23, 42), "text_muted": RGBColor(51, 65, 85), "tag_bg": RGBColor(209, 250, 229)},
        {"name": "Green",      "primary": RGBColor(22, 163, 74),  "secondary": RGBColor(34, 197, 94),   "accent": RGBColor(240, 253, 244), "card_bg": RGBColor(255, 255, 255), "border": RGBColor(187, 247, 208), "text_dark": RGBColor(15, 23, 42), "text_muted": RGBColor(51, 65, 85), "tag_bg": RGBColor(220, 252, 231)},
        {"name": "Indigo",     "primary": RGBColor(67, 56, 202),  "secondary": RGBColor(99, 102, 241),  "accent": RGBColor(238, 242, 255), "card_bg": RGBColor(255, 255, 255), "border": RGBColor(199, 210, 254), "text_dark": RGBColor(15, 23, 42), "text_muted": RGBColor(51, 65, 85), "tag_bg": RGBColor(224, 231, 255)},
        {"name": "Violet",     "primary": RGBColor(109, 40, 217), "secondary": RGBColor(139, 92, 246),  "accent": RGBColor(245, 243, 255), "card_bg": RGBColor(255, 255, 255), "border": RGBColor(221, 214, 254), "text_dark": RGBColor(15, 23, 42), "text_muted": RGBColor(51, 65, 85), "tag_bg": RGBColor(237, 233, 254)},
        {"name": "Purple",     "primary": RGBColor(126, 34, 206), "secondary": RGBColor(168, 85, 247),  "accent": RGBColor(250, 245, 255), "card_bg": RGBColor(255, 255, 255), "border": RGBColor(233, 213, 255), "text_dark": RGBColor(15, 23, 42), "text_muted": RGBColor(51, 65, 85), "tag_bg": RGBColor(243, 232, 255)},
        {"name": "Orange",     "primary": RGBColor(194, 65, 12),  "secondary": RGBColor(249, 115, 22),  "accent": RGBColor(255, 247, 237), "card_bg": RGBColor(255, 255, 255), "border": RGBColor(254, 215, 170), "text_dark": RGBColor(15, 23, 42), "text_muted": RGBColor(51, 65, 85), "tag_bg": RGBColor(255, 237, 213)},
        {"name": "Pink",       "primary": RGBColor(190, 24, 93),  "secondary": RGBColor(236, 72, 153),  "accent": RGBColor(253, 242, 248), "card_bg": RGBColor(255, 255, 255), "border": RGBColor(251, 207, 232), "text_dark": RGBColor(15, 23, 42), "text_muted": RGBColor(51, 65, 85), "tag_bg": RGBColor(252, 231, 243)},
        {"name": "Cyan",       "primary": RGBColor(14, 116, 144), "secondary": RGBColor(6, 182, 212),   "accent": RGBColor(236, 254, 255), "card_bg": RGBColor(255, 255, 255), "border": RGBColor(165, 243, 252), "text_dark": RGBColor(15, 23, 42), "text_muted": RGBColor(51, 65, 85), "tag_bg": RGBColor(207, 250, 254)}
    ]

    def add_base_decorations(slide, theme, slide_num, total_slides=11):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme["accent"]
        bg.line.color.rgb = theme["accent"]

        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = theme["primary"]
        top_bar.line.color.rgb = theme["primary"]

        footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.92), Inches(11.733), Inches(0.02))
        footer_bar.fill.solid()
        footer_bar.fill.fore_color.rgb = theme["border"]
        footer_bar.line.color.rgb = theme["border"]

        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.96), Inches(8.5), Inches(0.38))
        tf = footer_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "INTERNHUB  •  Full-Stack Internship & Recruitment Platform  •  NRCM / CSE-E"
        p.font.size = Pt(10)
        p.font.name = "Segoe UI"
        p.font.color.rgb = theme["text_muted"]

        slide_box = slide.shapes.add_textbox(Inches(10.5), Inches(6.96), Inches(2.033), Inches(0.38))
        stf = slide_box.text_frame
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.RIGHT
        sp.text = f"Slide {slide_num} of {total_slides}"
        sp.font.size = Pt(10.5)
        sp.font.bold = True
        sp.font.name = "Segoe UI"
        sp.font.color.rgb = theme["primary"]

    # ==========================================
    # SLIDE 1: Title & Student Details (Royal Blue)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    t1 = themes[0]
    add_base_decorations(slide1, t1, 1)

    # 1. Top Institution Header Banner
    inst_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.42), Inches(11.733), Inches(0.92))
    inst_card.fill.solid()
    inst_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    inst_card.line.color.rgb = t1["border"]
    inst_card.line.width = Pt(1.5)

    itf = inst_card.text_frame
    itf.word_wrap = True
    itp1 = itf.paragraphs[0]
    itp1.alignment = PP_ALIGN.CENTER
    itp1.text = "NARSIMHA REDDY ENGINEERING COLLEGE (NRCM)"
    itp1.font.size = Pt(16)
    itp1.font.bold = True
    itp1.font.name = "Segoe UI"
    itp1.font.color.rgb = t1["primary"]

    itp2 = itf.add_paragraph()
    itp2.alignment = PP_ALIGN.CENTER
    itp2.text = "DEPARTMENT OF COMPUTER SCIENCE & ENGINEERING  •  SUMMER INTERNSHIP PROJECT"
    itp2.font.size = Pt(11)
    itp2.font.bold = True
    itp2.font.name = "Segoe UI"
    itp2.font.color.rgb = t1["secondary"]
    itp2.space_before = Pt(3)

    # 2. Main Center Hero Project Banner (Compact & Sleek)
    hero_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.48), Inches(11.733), Inches(1.85))
    hero_card.fill.solid()
    hero_card.fill.fore_color.rgb = t1["primary"]
    hero_card.line.color.rgb = t1["secondary"]
    hero_card.line.width = Pt(1.5)

    htf = hero_card.text_frame
    htf.word_wrap = True
    
    hp1 = htf.paragraphs[0]
    hp1.alignment = PP_ALIGN.CENTER
    hp1.text = "INTERNHUB"
    hp1.font.size = Pt(30)
    hp1.font.bold = True
    hp1.font.name = "Segoe UI"
    hp1.font.color.rgb = RGBColor(255, 255, 255)

    hp2 = htf.add_paragraph()
    hp2.alignment = PP_ALIGN.CENTER
    hp2.text = "Internship Discovery & Recruitment Platform"
    hp2.font.size = Pt(15)
    hp2.font.bold = True
    hp2.font.name = "Segoe UI"
    hp2.font.color.rgb = RGBColor(219, 234, 254)
    hp2.space_before = Pt(3)

    hp3 = htf.add_paragraph()
    hp3.alignment = PP_ALIGN.CENTER
    hp3.text = "A Full-Stack Role-Based Platform Connecting Students, Recruiters & Administrators"
    hp3.font.size = Pt(11.5)
    hp3.font.name = "Segoe UI"
    hp3.font.color.rgb = RGBColor(191, 219, 254)
    hp3.space_before = Pt(4)

    hp4 = htf.add_paragraph()
    hp4.alignment = PP_ALIGN.CENTER
    hp4.text = "⚡ MongoDB  |  Express.js  |  React 18  |  Node.js  |  Docker  |  JWT & RBAC"
    hp4.font.size = Pt(11)
    hp4.font.bold = True
    hp4.font.name = "Segoe UI"
    hp4.font.color.rgb = RGBColor(254, 240, 138)
    hp4.space_before = Pt(6)

    # 3. Bottom Two Balanced Columns
    col_w = Inches(5.72)
    col_h = Inches(3.25)

    left_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.48), col_w, col_h)
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    left_card.line.color.rgb = t1["border"]
    left_card.line.width = Pt(1.5)

    ltf = left_card.text_frame
    ltf.word_wrap = True

    lp0 = ltf.paragraphs[0]
    lp0.text = "PRESENTED BY"
    lp0.font.size = Pt(11)
    lp0.font.bold = True
    lp0.font.name = "Segoe UI"
    lp0.font.color.rgb = t1["secondary"]

    lp1 = ltf.add_paragraph()
    lp1.text = "Candidate Details"
    lp1.font.size = Pt(16.5)
    lp1.font.bold = True
    lp1.font.name = "Segoe UI"
    lp1.font.color.rgb = t1["primary"]
    lp1.space_after = Pt(4)

    c_details = [
        ("Student Name", "RAATE NOMIKA"),
        ("Roll Number", "24X05A0533"),
        ("Branch", "B.Tech — Computer Science & Engineering"),
        ("Section", "CSE-E"),
        ("Academic Year", "2023 – 2027")
    ]

    for label, val in c_details:
        p = ltf.add_paragraph()
        p.text = f"•  {label}: "
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.font.color.rgb = t1["primary"]
        p.space_after = Pt(2)
        
        run = p.add_run()
        run.text = val
        run.font.bold = False
        run.font.color.rgb = t1["text_dark"]

    right_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.81), Inches(3.48), col_w, col_h)
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    right_card.line.color.rgb = t1["border"]
    right_card.line.width = Pt(1.5)

    rtf = right_card.text_frame
    rtf.word_wrap = True

    rp0 = rtf.paragraphs[0]
    rp0.text = "PROJECT SPECIFICATION"
    rp0.font.size = Pt(11)
    rp0.font.bold = True
    rp0.font.name = "Segoe UI"
    rp0.font.color.rgb = t1["secondary"]

    rp1 = rtf.add_paragraph()
    rp1.text = "Academic & Technical Scope"
    rp1.font.size = Pt(16.5)
    rp1.font.bold = True
    rp1.font.name = "Segoe UI"
    rp1.font.color.rgb = t1["primary"]
    rp1.space_after = Pt(4)

    p_details = [
        ("Project Domain", "Full-Stack Web Development & Cloud Computing"),
        ("Architecture", "MERN 3-Tier Decoupled Client-Server Model"),
        ("Core Capability", "Role-Based Access Control (Student / Recruiter / Admin)"),
        ("Deployment", "Dockerized Containerization & Netlify Cloud Hosting"),
        ("Testing", "Automated E2E & API Integration Test Lifecycle")
    ]

    for label, val in p_details:
        p = rtf.add_paragraph()
        p.text = f"•  {label}: "
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.font.color.rgb = t1["primary"]
        p.space_after = Pt(2)
        
        run = p.add_run()
        run.text = val
        run.font.bold = False
        run.font.color.rgb = t1["text_dark"]

    slide1.notes_slide.notes_text_frame.text = (
        "Good morning / afternoon respected professors, evaluators, and my dear peers.\n"
        "My name is RAATE NOMIKA, bearing roll number 24X05A0533, Section CSE-E from the Department of Computer Science & Engineering "
        "at Narsimha Reddy Engineering College (NRCM).\n"
        "Today, I am proud to present our project: INTERNHUB — a Full-Stack MERN Internship Discovery & Recruitment Platform."
    )

    # Master Content Slide Generator (Zero White Space)
    def add_content_slide(slide_index, title, subtitle, items, note_text, layout_type="two_column", takeaway_text=""):
        th = themes[slide_index]
        slide = prs.slides.add_slide(blank_layout)
        add_base_decorations(slide, th, slide_index + 1)

        # Header Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.42), Inches(1.3), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = th["tag_bg"]
        badge.line.color.rgb = th["border"]
        btf = badge.text_frame
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        bp.text = f"SLIDE {slide_index + 1:02d}"
        bp.font.size = Pt(10.5)
        bp.font.bold = True
        bp.font.name = "Segoe UI"
        bp.font.color.rgb = th["primary"]

        # Title Box
        tbox = slide.shapes.add_textbox(Inches(2.25), Inches(0.32), Inches(10.2), Inches(0.85))
        ttframe = tbox.text_frame
        ttframe.word_wrap = True
        tp = ttframe.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(22)
        tp.font.bold = True
        tp.font.name = "Segoe UI"
        tp.font.color.rgb = th["primary"]

        tsp = ttframe.add_paragraph()
        tsp.text = subtitle
        tsp.font.size = Pt(12)
        tsp.font.name = "Segoe UI"
        tsp.font.color.rgb = th["text_muted"]
        tsp.space_before = Pt(2)

        # 4-Grid Layout
        if layout_type == "grid_4":
            card_w = Inches(5.72)
            card_h = Inches(2.28)
            coords = [
                (Inches(0.8), Inches(1.32)),
                (Inches(6.81), Inches(1.32)),
                (Inches(0.8), Inches(3.72)),
                (Inches(6.81), Inches(3.72))
            ]

            for i, (item_title, item_desc) in enumerate(items[:4]):
                x, y = coords[i]
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
                card.fill.solid()
                card.fill.fore_color.rgb = th["card_bg"]
                card.line.color.rgb = th["border"]
                card.line.width = Pt(1.5)

                # Number pill icon
                pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.28), y + Inches(0.22), Inches(0.58), Inches(0.38))
                pill.fill.solid()
                pill.fill.fore_color.rgb = th["tag_bg"]
                pill.line.color.rgb = th["border"]
                ptf = pill.text_frame
                pp = ptf.paragraphs[0]
                pp.alignment = PP_ALIGN.CENTER
                pp.text = f"{i+1:02d}"
                pp.font.size = Pt(12)
                pp.font.bold = True
                pp.font.name = "Segoe UI"
                pp.font.color.rgb = th["primary"]

                # Card content textbox
                cbox = slide.shapes.add_textbox(x + Inches(1.02), y + Inches(0.16), card_w - Inches(1.25), card_h - Inches(0.25))
                ctf = cbox.text_frame
                ctf.word_wrap = True

                ip = ctf.paragraphs[0]
                ip.text = item_title
                ip.font.size = Pt(14.5)
                ip.font.bold = True
                ip.font.name = "Segoe UI"
                ip.font.color.rgb = th["primary"]

                dp = ctf.add_paragraph()
                dp.text = item_desc
                dp.font.size = Pt(12.5)
                dp.font.name = "Segoe UI"
                dp.font.color.rgb = th["text_muted"]
                dp.space_before = Pt(6)

        # 2-Column Layout (With perfectly distributed blocks to eliminate empty bottom space)
        elif layout_type == "two_column":
            col_width = Inches(5.72)
            left_pos = [Inches(0.8), Inches(6.81)]
            mid = (len(items) + 1) // 2
            col_items = [items[:mid], items[mid:]]

            for i in range(2):
                sec_title, point_list = col_items[i][0]

                # 1. Column Header Badge Card
                hdr_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos[i], Inches(1.30), col_width, Inches(0.68))
                hdr_card.fill.solid()
                hdr_card.fill.fore_color.rgb = th["primary"]
                hdr_card.line.color.rgb = th["secondary"]
                hdr_card.line.width = Pt(1.2)
                
                htf = hdr_card.text_frame
                htf.word_wrap = True
                hp = htf.paragraphs[0]
                hp.alignment = PP_ALIGN.CENTER
                hp.text = sec_title
                hp.font.size = Pt(15.5)
                hp.font.bold = True
                hp.font.name = "Segoe UI"
                hp.font.color.rgb = RGBColor(255, 255, 255)

                # 2. Individual Content Cards for each point (fills the column completely)
                num_pts = len(point_list)
                pt_card_h = Inches(0.92)
                pt_gap = Inches(0.09)
                start_y = Inches(2.06)

                for p_idx, (bp_title, bp_desc) in enumerate(point_list):
                    cur_y = start_y + p_idx * (pt_card_h + pt_gap)
                    p_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos[i], cur_y, col_width, pt_card_h)
                    p_card.fill.solid()
                    p_card.fill.fore_color.rgb = th["card_bg"]
                    p_card.line.color.rgb = th["border"]
                    p_card.line.width = Pt(1.2)

                    pctf = p_card.text_frame
                    pctf.word_wrap = True
                    
                    # Point Title
                    pp0 = pctf.paragraphs[0]
                    pp0.text = f"✔  {bp_title}"
                    pp0.font.size = Pt(13)
                    pp0.font.bold = True
                    pp0.font.name = "Segoe UI"
                    pp0.font.color.rgb = th["primary"]

                    # Point Description
                    pp1 = pctf.add_paragraph()
                    pp1.text = bp_desc
                    pp1.font.size = Pt(11.5)
                    pp1.font.name = "Segoe UI"
                    pp1.font.color.rgb = th["text_muted"]
                    pp1.space_before = Pt(2)

        # Bottom Key Highlight Bar (Anchors the slide)
        if takeaway_text:
            t_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.15), Inches(11.733), Inches(0.62))
            t_bar.fill.solid()
            t_bar.fill.fore_color.rgb = th["tag_bg"]
            t_bar.line.color.rgb = th["border"]
            t_bar.line.width = Pt(1.2)

            t_tf = t_bar.text_frame
            t_tf.word_wrap = True
            tp = t_tf.paragraphs[0]
            tp.text = f"💡 Impact Summary: {takeaway_text}"
            tp.font.size = Pt(12)
            tp.font.bold = True
            tp.font.name = "Segoe UI"
            tp.font.color.rgb = th["primary"]

        # Speaker notes
        slide.notes_slide.notes_text_frame.text = note_text

    # ==========================================
    # SLIDE 2: Problem Statement & Motivation (Sky Blue)
    # ==========================================
    add_content_slide(
        1,
        "Problem Statement & Motivation",
        "Key challenges in existing student recruitment methods and why InternHub was built",
        [
            ("Fragmented Discovery Channels", "Students have to navigate multiple generic job boards filled with experienced roles, resulting in wasted effort and missed entry-level opportunities."),
            ("Application Black Hole", "Lack of transparent tracking leaves candidates with zero feedback on whether resumes were seen, reviewed, shortlisted, or rejected."),
            ("Recruiter Screening Overhead", "Employers receive disorganized emails and unstandardized resumes, making manual candidate screening slow, tedious, and prone to errors."),
            ("Lack of a Unified Trusted Hub", "Absence of a single role-based platform bridging students, corporate recruiters, and college administrators seamlessly.")
        ],
        '"Respected professors and peers, every year millions of college students struggle to find verified internship opportunities. Most job portals are cluttered with senior roles, and students face a \'black hole\' where they never know whether their resume was reviewed, shortlisted, or rejected. On the other hand, employers spend hours manually sorting through unorganized emails and resumes. To eliminate this gap, we built InternHub—a dedicated, transparent, and structured platform tailored specifically for college students and recruiters."',
        "grid_4",
        "Solves recruitment latency and transparency gaps with real-time tracking and verified opportunities."
    )

    # ==========================================
    # SLIDE 3: Project Overview & Core Objectives (Teal)
    # ==========================================
    add_content_slide(
        2,
        "Project Overview & Core Objectives",
        "Strategic purpose, technical mission, and key deliverables of the InternHub portal",
        [
            ("Modern Full-Stack MERN Architecture", "Engineered with MongoDB, Express.js, React 18, and Node.js for lightning-fast responsiveness and full-stack modularity."),
            ("Smart Opportunity Discovery", "Multi-parameter search allowing candidates to instantly filter openings by domain, location (remote/hybrid), duration, and stipend."),
            ("Transparent Live Tracking", "End-to-end status visibility across every stage: Applied ➔ In Review ➔ Shortlisted ➔ Decision, eliminating student uncertainty."),
            ("Recruiter & Admin Dashboards", "Comprehensive portals for employers to manage applicants and for administrators to oversee platform moderation and metrics.")
        ],
        '"InternHub is an end-to-end web portal built on the modern MERN stack. Our primary objective is to make the hiring workflow effortless. Students can discover relevant opportunities with advanced filters and track their application status live. Employers can post openings, review applicants, and update hiring statuses with a single click. Meanwhile, the administrator ensures platform integrity and data authenticity through an admin dashboard."',
        "grid_4",
        "Delivers a zero-latency recruitment pipeline tailored for university students and corporate recruiters."
    )

    # ==========================================
    # SLIDE 4: System Architecture & Technology Stack (Emerald)
    # ==========================================
    add_content_slide(
        3,
        "System Architecture & Technology Stack",
        "Decoupled Client-Server Tiered Architecture with Scalable RESTful Services",
        [
            ("Frontend & Client Tier", [
                ("React 18 & Vite", "High-performance reactive frontend with ultra-fast Hot Module Replacement"),
                ("TailwindCSS & Modern UI", "Responsive design system delivering cross-device accessibility and animations"),
                ("React Router v6", "Dynamic client-side routing with authenticated route guards and error boundaries"),
                ("Axios HTTP Client", "Centralized API communication layer with request/response interceptors")
            ]),
            ("Backend & Persistence Tier", [
                ("Node.js & Express.js", "Asynchronous, event-driven RESTful API backend handling concurrent requests"),
                ("MongoDB & Mongoose ODM", "Document-oriented database with strict schema validation and query indexing"),
                ("JWT Authentication", "Stateless JSON Web Tokens with encrypted Bearer header authorization"),
                ("Docker & Cloud Deploy", "Containerized deployment ensuring environment parity across local and cloud")
            ])
        ],
        '"Coming to the technical architecture, our system follows a decoupled Client-Server model. On the frontend, we use React 18 bundled with Vite for fast performance, responsive state management, and a clean user experience. The backend is powered by Node.js and Express.js, providing a scalable REST API. We use MongoDB as our database with Mongoose ODM to handle dynamic documents like resumes, applications, and user profiles."',
        "two_column",
        "Decoupled MERN stack guarantees high throughput, modularity, and rapid cloud deployment."
    )

    # ==========================================
    # SLIDE 5: Role-Based Access Control (RBAC) (Green)
    # ==========================================
    add_content_slide(
        4,
        "Role-Based Access Control (RBAC)",
        "Granular access control and dedicated dashboards for Students, Employers, and Admins",
        [
            ("Student / Candidate Role", [
                ("Opportunity Discovery", "Explore verified internships with smart domain, location, and stipend filters"),
                ("One-Click Apply", "Submit applications with linked profile details, resume URL, and cover notes"),
                ("Live Status Tracking", "Real-time visibility into application stage (Applied, Review, Shortlisted)"),
                ("Bookmark & Alerts", "Save favorite opportunities and receive instant status update notifications")
            ]),
            ("Employer & Administrator Roles", [
                ("Job Management", "Recruiters can create, edit, publish, and close internship opportunities"),
                ("Applicant Evaluation Board", "Review profiles, download resumes, and toggle hiring status with 1-click"),
                ("Platform Moderation", "Admins verify recruiters, audit job postings, and manage user permissions"),
                ("System Analytics", "Oversee total placements, active postings, and user engagement metrics")
            ])
        ],
        '"Security and user segmentation are critical in our design. We implemented Role-Based Access Control, dividing users into three distinct roles: Students, Employers, and Admins. Each role has protected routes and dedicated dashboards. A student cannot access employer tools, and an employer cannot alter administrative settings. This ensures data segregation and a clutter-free interface tailored to each user\'s specific tasks."',
        "two_column",
        "Strict data isolation and customized workflows for students, corporate recruiters, and college admins."
    )

    # ==========================================
    # SLIDE 6: Key Features & Functional Modules (Indigo)
    # ==========================================
    add_content_slide(
        5,
        "Key Features & Functional Modules",
        "Core functional workflows powering user interaction and platform efficiency",
        [
            ("Smart Search & Filter Engine", "Multi-parameter query system supporting domain keywords, location (remote/hybrid/office), stipend range, and job duration."),
            ("One-Click Application System", "Instant application submission linking uploaded student resume, GitHub/portfolio links, and customized cover notes."),
            ("Real-Time Application Pipeline", "Transparent status transitions: Applied ➔ In Review ➔ Shortlisted ➔ Decision with timestamped updates."),
            ("Recruiter Decision Workspace", "Kanban/Table style candidate evaluation board allowing recruiters to review profiles, preview resumes, and toggle hiring status.")
        ],
        '"Let\'s look at the core functional modules. First is our Smart Search and Filter engine, allowing students to quickly filter roles by location, stipend, skill set, and job type. Second is the one-click application submission where students attach their profiles and resumes. Third is the real-time application pipeline, giving candidates full transparency on whether their profile is in review or shortlisted. For employers, we built an intuitive candidate evaluation board to streamline decisions."',
        "grid_4",
        "Streamlines the hiring process from candidate discovery to final offer confirmation."
    )

    # ==========================================
    # SLIDE 7: Database Design & Data Modeling (Violet)
    # ==========================================
    add_content_slide(
        6,
        "Database Design & Data Modeling",
        "Optimized MongoDB Collections with relational references and compound indexing",
        [
            ("Core Database Collections", [
                ("Users Collection", "Stores credentials, hashed passwords, user roles (Student/Employer/Admin), and profiles"),
                ("Internships Collection", "Job title, company, recruiter reference, requirements, stipend, location, and deadlines"),
                ("Applications Collection", "Relational junction linking Candidate ID, Internship ID, resume URL, and status enum"),
                ("Notifications Collection", "System alerts, status update logs, timestamps, and user read-receipt tracking")
            ]),
            ("Integrity & Query Optimization", [
                ("Relational Referencing", "Normalized links via Mongoose ObjectId population for clean relational queries"),
                ("Compound Indexing", "Indexed search keys on skills, location, and role achieving sub-50ms response times"),
                ("Schema Validation", "Strict enum state validation preventing invalid status transitions and corrupted records"),
                ("Cascading Cleanups", "Automated reference cleanup on internship closure or account termination")
            ])
        ],
        '"Our database model in MongoDB is designed for high read throughput and relational integrity. We maintain collections for Users, Internships, and Applications. The Applications schema acts as a junction linking the Student ID to the Internship ID with status states such as \'Applied\', \'Under Review\', \'Shortlisted\', and \'Rejected\'. We utilized database indexing on frequent search keys like location and skills to ensure fast response times even as data grows."',
        "two_column",
        "MongoDB schema architecture delivers sub-second search speeds and reliable transactional data."
    )

    # ==========================================
    # SLIDE 8: Security, Authentication & Best Practices (Purple)
    # ==========================================
    add_content_slide(
        7,
        "Security, Authentication & Best Practices",
        "Multi-layered defense architecture securing credentials, sessions, and APIs",
        [
            ("Authentication & Session Security", [
                ("Bcrypt Password Hashing", "Adaptive salt rounds password encryption preventing plaintext credential leaks"),
                ("Stateless JWT Tokens", "Digitally signed tokens verified on every protected API call via Bearer headers"),
                ("Frontend Route Guards", "Automatic redirection and token verification preventing unauthorized URL visits"),
                ("Role Verification Middleware", "Server-level RBAC checks ensuring students cannot invoke recruiter endpoints")
            ]),
            ("Server Hardening & Defenses", [
                ("Helmet.js Integration", "Sets robust HTTP security headers to protect against clickjacking and XSS attacks"),
                ("Express Rate Limiting", "Throttles repeated API requests to shield login endpoints from brute-force attacks"),
                ("Strict CORS Policies", "Restricts cross-origin resource sharing strictly to verified client domain origins"),
                ("Input Sanitization", "Validates and sanitizes request payloads to prevent NoSQL injection vulnerabilities")
            ])
        ],
        '"Security is implemented at every layer of InternHub. User passwords are encrypted using bcrypt hashing before storage. Authentication is handled via stateless JSON Web Tokens (JWT) verified on every protected API call. To safeguard the server, we added security headers using Helmet, strict CORS policies, and rate limiters to protect against brute-force login attempts and DDoS attacks."',
        "two_column",
        "Enterprise-grade security stack shielding against injection, session hijacking, and brute-force."
    )

    # ==========================================
    # SLIDE 9: UI/UX & Frontend Engineering (Orange)
    # ==========================================
    add_content_slide(
        8,
        "UI/UX & Frontend Engineering",
        "Crafting a responsive, accessible, and high-performance user interface",
        [
            ("Component-Driven Architecture", "Modular React components (Cards, Modals, Status Badges, Toast Alerts, Tables) ensuring UI consistency and easy maintainability across all pages."),
            ("Cross-Device Responsiveness", "Fluid layout design optimized for mobile smartphones, tablets, and wide-screen desktop displays using TailwindCSS breakpoints."),
            ("Interactive State Management", "Instant UI feedback upon applying or updating candidate statuses without requiring full-page browser reloads."),
            ("SEO & Fast Asset Loading", "Dynamic OpenGraph and Meta tags via SEOHead component, lazy-loaded routes, and lightweight SVG icons (Lucide React).")
        ],
        '"For the frontend, we focused heavily on UX and responsiveness. The portal adapts smoothly across mobile devices, tablets, and desktops. Using React\'s modular component structure, we created reusable cards, modals, and tables. We also integrated client-side route guards so unauthenticated users are redirected automatically. Additionally, we implemented SEO best practices and fast asset loading for optimal performance."',
        "grid_4",
        "Rich, accessible user interface optimized for high engagement and cross-device speed."
    )

    # ==========================================
    # SLIDE 10: Testing, Containerization & Deployment (Pink)
    # ==========================================
    add_content_slide(
        9,
        "Testing, Containerization & Deployment",
        "Ensuring software reliability, reproducible environments, and continuous delivery",
        [
            ("Automated Testing & QA", [
                ("Integration Testing", "Validates complete request-response lifecycles across all authentication & job APIs"),
                ("End-to-End (E2E) Suites", "Simulates realistic student journeys: Registration ➔ Job Search ➔ Application"),
                ("Security Boundary Checks", "Tests expired tokens, unauthorized role escalations, and invalid payload handling"),
                ("Error Boundary Wrappers", "Prevents entire React client crashes by catching and displaying graceful UI fallbacks")
            ]),
            ("DevOps & Production Hosting", [
                ("Docker & Docker Compose", "Containerized microservices guaranteeing identical environments on dev and prod"),
                ("Frontend Deployment", "Hosted on Netlify with automated continuous deployment and global CDN delivery"),
                ("Backend Cloud Runtimes", "Express server deployed with isolated environment variables and auto-restarts"),
                ("Database Cloud Cluster", "MongoDB Atlas replica set featuring automated daily backups and live monitoring")
            ])
        ],
        '"To ensure reliability, we wrote comprehensive integration and end-to-end test suites verifying the entire application lifecycle—from user registration and job posting to application submission. We also containerized the application using Docker and Docker Compose, ensuring zero environment discrepancies between development and production. The live frontend is deployed on Netlify, connected to our cloud-hosted backend."',
        "two_column",
        "Dockerized architecture paired with automated testing ensures zero-downtime production deployment."
    )

    # ==========================================
    # SLIDE 11: Future Scope & Conclusion (Cyan)
    # ==========================================
    add_content_slide(
        10,
        "Future Scope & Project Conclusion",
        "Strategic innovation roadmap and final summary of project achievements",
        [
            ("AI-Powered Candidate Match", "Machine learning algorithms to automatically parse student resumes, calculate skill match scores, and recommend optimal internships."),
            ("In-App Real-Time Communication", "Integrated WebRTC chat and virtual interview scheduling directly within the recruiter evaluation board."),
            ("Cross-Platform Mobile App", "React Native mobile companion app for iOS and Android with push notifications for instant status changes."),
            ("Project Conclusion", "InternHub successfully bridges the gap between students and recruiters through a modern, secure, and transparent recruitment platform.")
        ],
        '"To conclude, InternHub solves a real-world problem by providing a modern, transparent, and secure hiring bridge between students and recruiters. In future iterations, we plan to integrate AI-driven resume scoring to match candidates automatically with relevant jobs, real-time messaging, and direct integrations with college placement cells. Thank you, professors. We are now open to any questions."',
        "grid_4",
        "InternHub modernizes university hiring, providing a scalable blueprint for campus placements."
    )

    output_path = r"c:\internship\InternHub_College_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_internhub_presentation()
